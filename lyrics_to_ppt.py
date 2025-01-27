from pptx import Presentation
from pptx.util import Inches

def create_lyrics_presentation(lyrics, output_file):
    # Create a presentation object
    prs = Presentation()

    # Define slide layout (blank layout)
    blank_slide_layout = prs.slide_layouts[6]

    # Split lyrics into lines
    lines = lyrics.strip().split('\n')

    # Process lyrics in chunks of 2 lines per slide
    for i in range(0, len(lines), 2):
        # Get 2 lines for the current slide
        slide_lines = lines[i:i+2]
        slide_text = "\n".join(slide_lines)

        # Add a new slide
        slide = prs.slides.add_slide(blank_slide_layout)

        # Set background color to green
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = (0, 128, 0)  # RGB for green

        # Add a text box for the lyrics
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(4)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame

        # Add lyrics to the text box
        p = text_frame.add_paragraph()
        p.text = slide_text
        p.font.size = 40
        p.font.color.rgb = (255, 255, 255)  # White text
        p.alignment = 1  # Center alignment

    # Save the presentation
    prs.save(output_file)

# Example usage
lyrics = """
Here's to the ones that we got
Cheers to the wish you were here, but you're not
'Cause the drinks bring back all the memories
Of everything we've been through
Toast to the ones here today
Toast to the ones that we lost on the way
'Cause the drinks bring back all the memories
And the memories bring back, memories bring back you
"""

create_lyrics_presentation(lyrics, "lyrics_presentation.pptx")
print("Presentation created successfully!")
